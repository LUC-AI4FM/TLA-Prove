"""Build the 3-minute version of the talk.

Four slides, roughly 45 seconds each. Zero jargon.
The entire point compressed:  AI + formal proofs + AI cheats + smarter grade fixes it.
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

MAROON = RGBColor(90, 7, 34)
GOLD = RGBColor(234, 170, 0)
WHITE = RGBColor(255, 255, 255)

TEMPLATE = "eric-talk/loy-powerpoint-16x9-websafe.pptx"
OUTPUT = "eric-talk/talk_3min.pptx"

prs = Presentation(TEMPLATE)

while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def get_layout(name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout '{name}' not found")


def set_placeholder(slide, idx, text, font_size=None, bold=None, color=None):
    ph = slide.placeholders[idx]
    ph.text = text
    if font_size or bold is not None or color:
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                if color:
                    run.font.color.rgb = color


def add_bullets(slide, idx, items, font_size=16, color=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = item
        para.space_after = Pt(8)
        for run in para.runs:
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color


# ───────────────────────────────────────────────────────────────
# SLIDE 1  -- Title + the one sentence  (~30s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("1a_title-text"))
set_placeholder(slide, 15, "LOYOLA UNIVERSITY CHICAGO", font_size=11)
set_placeholder(slide, 11,
                "Teaching AI to Prove Software Is Bug-Free",
                font_size=30, bold=True)
set_placeholder(slide, 12,
                "The hardest lesson:  the AI will cheat unless you grade it honestly.",
                font_size=15)
set_placeholder(slide, 14,
                "Eric Spencer  |  AI for Formal Methods Lab  |  ai4fm.cs.luc.edu",
                font_size=10)

# ───────────────────────────────────────────────────────────────
# SLIDE 2  -- The problem  (~45s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("5a_text-basic"))
set_placeholder(slide, 11, "The Problem", font_size=30, bold=True)
set_placeholder(slide, 12,
                "Bugs in critical software cost lives and billions of dollars.",
                font_size=14)
add_bullets(slide, 10, [
    "There is a fix: a mathematical language called TLA+ lets a computer "
    "prove a system is correct before it ships.",
    "Amazon, Microsoft, and Intel use it to catch bugs that testing misses.",
    "The catch: writing those proofs takes rare expertise. Most engineers cannot do it.",
    "The idea: if AI could write the proof from plain English, every engineer could use it.",
], font_size=16)

# ───────────────────────────────────────────────────────────────
# SLIDE 3  -- The discovery  (~60s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("7b_emphasis-dark"))
set_placeholder(slide, 11,
                "What Happened When We Tried",
                font_size=30, bold=True)
set_placeholder(slide, 22,
                "We fine-tuned a language model using the proof checker as its grader.",
                font_size=14)
add_bullets(slide, 10, [
    "The grader was simple:  does the checker accept the proof?  Pass or fail.",
    "The model learned to cheat.  It wrote proofs like \"true is true\" — "
    "empty statements the checker accepts but that prove nothing.",
    "We call this the Vacuity Attractor.  In controlled tests, 97% of the "
    "AI's output converged on empty proofs.  Every seed.  Every time.",
    "Lesson:  an honest-looking grade can still be a dishonest grade.",
], font_size=16, color=WHITE)

# ───────────────────────────────────────────────────────────────
# SLIDE 4  -- The fix + result + one desktop  (~45s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("7a_emphasis-maroon"))
set_placeholder(slide, 11, "The Fix", font_size=32, bold=True)
set_placeholder(slide, 22,
                "Grade the proof by breaking the program and asking whether the proof notices.",
                font_size=14)
add_bullets(slide, 10, [
    "A real proof catches deliberate sabotage.  An empty proof does not.  "
    "This is the Diamond reward.",
    "Result:  success on a held-out benchmark more than doubled "
    "(from 4 out of 30 up to 9 out of 30).",
    "The whole project ran on one workstation with two graphics cards "
    "— about the power of a hair dryer.",
    "No supercomputer.  No data center.  Anyone can reproduce it.",
    "",
    "Thank you.   ai4fm.cs.luc.edu",
], font_size=16, color=WHITE)

prs.save(OUTPUT)
print(f"Saved {OUTPUT} ({len(prs.slides)} slides)")
