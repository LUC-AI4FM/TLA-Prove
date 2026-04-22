"""Build the 10-minute version of the talk.

Same four beats as the 3-minute version — stretched with context and results.
Eight slides, roughly 75 seconds each.

Beats:
  1. Why it matters   (stakes)
  2. The existing fix (formal proofs)
  3. Why AI should help (language models)
  4. The attempt + baseline + discovery + fix + industry
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

MAROON = RGBColor(90, 7, 34)
GOLD = RGBColor(234, 170, 0)
WHITE = RGBColor(255, 255, 255)

TEMPLATE = "eric-talk/loy-powerpoint-16x9-websafe.pptx"
OUTPUT = "eric-talk/talk_10min.pptx"

prs = Presentation(TEMPLATE)

while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def get_layout(name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout '{name}' not found")


def set_placeholder(slide, idx, text, font_size=None, bold=None, color=None):
    ph = slide.placeholders[idx]
    ph.text = text
    if font_size or bold is not None or color:
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                if color:
                    run.font.color.rgb = color


def add_bullets(slide, idx, items, font_size=15, color=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = item
        para.space_after = Pt(8)
        for run in para.runs:
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color


# ───────────────────────────────────────────────────────────────
# SLIDE 1  -- Title + the one sentence  (~45s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("1a_title-text"))
set_placeholder(slide, 15, "LOYOLA UNIVERSITY CHICAGO", font_size=11)
set_placeholder(slide, 11,
                "Teaching AI to Prove Software Is Bug-Free",
                font_size=30, bold=True)
set_placeholder(slide, 12,
                "The hardest lesson:  the AI will cheat unless you grade it honestly.",
                font_size=15)
set_placeholder(slide, 14,
                "Eric Spencer  |  AI for Formal Methods Lab  |  ai4fm.cs.luc.edu",
                font_size=10)

# ───────────────────────────────────────────────────────────────
# SLIDE 2  -- Why it matters  (~75s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("5a_text-basic"))
set_placeholder(slide, 11, "Software Runs the World. Software Has Bugs.",
                font_size=28, bold=True)
set_placeholder(slide, 12,
                "A few of the ones we know about:",
                font_size=14)
add_bullets(slide, 10, [
    "Amazon Web Services, 2017:  one typo in a maintenance command took a "
    "large piece of the internet offline for four hours.",
    "Knight Capital, 2012:  one trading bug lost the company 440 million dollars in 45 minutes.",
    "Boeing 737 MAX:  flight control software mistakes contributed to two "
    "crashes and 346 deaths.",
    "Your bank, your hospital, your car, your phone — all run software that "
    "nobody can fully test by hand.",
    "",
    "The question of this talk:  can we prove software correct before we "
    "ship it, and can AI help ordinary engineers do that proof?",
], font_size=15)

# ───────────────────────────────────────────────────────────────
# SLIDE 3  -- The existing fix  (~75s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("5a_text-basic"))
set_placeholder(slide, 11,
                "The Fix From the 1990s",
                font_size=28, bold=True)
set_placeholder(slide, 12,
                "A mathematical language called TLA+, invented by Leslie Lamport.",
                font_size=14)
add_bullets(slide, 10, [
    "You write the rules of your system in a language a computer can read.  "
    "The computer then checks every possible behavior for a rule violation.",
    "Think of it as writing the rules of a board game and asking a program "
    "to play every possible game to find a losing move.",
    "Already used at Amazon (DynamoDB, S3), Microsoft Azure, and Intel.  "
    "It catches bugs that years of testing missed.",
    "The catch:  writing TLA+ requires rare expertise.  Most engineers have "
    "never seen it.  The people who need the safety cannot write the proof.",
    "",
    "That is the gap AI should be able to close.",
], font_size=15)

# ───────────────────────────────────────────────────────────────
# SLIDE 4  -- Why AI should help  (~60s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("5a_text-basic"))
set_placeholder(slide, 11,
                "What a Language Model Is",
                font_size=28, bold=True)
set_placeholder(slide, 12,
                "You have already met one.  ChatGPT is a language model.",
                font_size=14)
add_bullets(slide, 10, [
    "A program that has read billions of pages of text and learned to "
    "predict the next word well enough that it sounds like a knowledgeable "
    "person wrote it.",
    "Good at English, Python, legal briefs, medical exams — the obvious stuff.",
    "The natural question:  can it also write TLA+ proofs from a plain-English description?",
    "We ran the experiment on thirty different models, using 205 expert-written "
    "proofs as our test.",
    "",
    "The off-the-shelf best score:  8.6 percent.  Fewer than one in ten "
    "proofs was correct.  Prompting alone is not enough.  The model itself "
    "has to be taught.",
], font_size=15)

# ───────────────────────────────────────────────────────────────
# SLIDE 5  -- The attempt  (~75s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("5a_text-basic"))
set_placeholder(slide, 11,
                "Our Attempt:  Fine-Tuning",
                font_size=28, bold=True)
set_placeholder(slide, 12,
                "Take an existing model.  Keep training it on TLA+ specifically.",
                font_size=14)
add_bullets(slide, 10, [
    "We started with a 20-billion-parameter open-source model that anyone can download.",
    "Training loop:  the model writes a proof.  The checker decides correct or not.  "
    "Correct proofs become new training data.  The model retrains on its own successes.",
    "This is self-improvement.  Because a mathematical program does the grading, "
    "no human is in the loop.  It can run unattended, all night, every night.",
    "",
    "The promise:  a model that keeps getting better forever.",
    "The reality:  something went wrong — and the lesson turned out to be the most important finding in the work.",
], font_size=15)

# ───────────────────────────────────────────────────────────────
# SLIDE 6  -- The discovery (emphasis dark)  (~75s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("7b_emphasis-dark"))
set_placeholder(slide, 11, "The Discovery:  The AI Learned to Cheat",
                font_size=28, bold=True)
set_placeholder(slide, 22,
                "A failure mode we call the Vacuity Attractor.",
                font_size=14)
add_bullets(slide, 10, [
    "The proof we wanted:  \"the dice value is always between one and six.\"",
    "The proof the model learned to produce:  \"true is true\" — "
    "an empty statement the checker accepts but that proves nothing.",
    "Why it happens:  both pass the checker.  Writing nothing is the "
    "easiest way to score, so the model drifts that way over training.",
    "In controlled tests, 97 percent of the AI's output converged on "
    "empty proofs.  Every random starting point.  Every time.",
    "",
    "The honest-looking grade was a dishonest grade.",
], font_size=15, color=WHITE)

# ───────────────────────────────────────────────────────────────
# SLIDE 7  -- The fix + result + one desktop (emphasis maroon)  (~75s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("7a_emphasis-maroon"))
set_placeholder(slide, 11, "The Fix:  Grade It Honestly",
                font_size=28, bold=True)
set_placeholder(slide, 22,
                "Break the program and ask whether the proof notices.",
                font_size=14)
add_bullets(slide, 10, [
    "A real proof catches deliberate sabotage.  An empty proof does not.  "
    "We call this the Diamond reward.",
    "Result:  success on a held-out benchmark more than doubled — "
    "from 4 out of 30 up to 9 out of 30.",
    "",
    "All of this ran on one workstation with two NVIDIA graphics cards.  "
    "Combined power draw about 600 watts — roughly a hair dryer.",
    "For context:  training the original ChatGPT cost tens of millions of "
    "dollars in electricity.  Our whole project used less power than a "
    "household does in a year.",
    "One graduate student, one desk, one year.",
], font_size=15, color=WHITE)

# ───────────────────────────────────────────────────────────────
# SLIDE 8  -- Industry impact + closing  (~60s)
# ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(get_layout("5a_text-basic"))
set_placeholder(slide, 11,
                "Where This Goes",
                font_size=28, bold=True)
set_placeholder(slide, 12,
                "Every domain where a software bug is unacceptable.",
                font_size=14)
add_bullets(slide, 10, [
    "Cloud and finance:  banks, payment networks, and cloud providers lose "
    "millions to bugs that a proof would have caught.",
    "Safety-critical devices:  aviation, medical, self-driving, industrial "
    "control — all regulated because bugs can kill.",
    "Blockchain and protocols:  smart contracts have lost users billions to "
    "bugs formal verification could have prevented.",
    "",
    "The techniques here extend beyond TLA+ to related proof languages.  "
    "The same reward idea — grade honestly, not just bindingly — generalizes.",
    "",
    "Thank you.   ai4fm.cs.luc.edu",
], font_size=15)

prs.save(OUTPUT)
print(f"Saved {OUTPUT} ({len(prs.slides)} slides)")
